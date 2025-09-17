import os
import logging
from typing import Dict, List, Any, Optional, TypedDict, Annotated
from datetime import datetime
from pathlib import Path
import tempfile
from io import BytesIO
import json
import io

from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse , StreamingResponse
from pydantic import BaseModel
import uvicorn

# LangGraph imports
from langgraph.graph import StateGraph, END
from langgraph.graph.message import add_messages
from langchain_core.messages import BaseMessage

# LangChain imports
from langchain_groq import ChatGroq
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEndpointEmbeddings
from langchain_core.output_parsers import JsonOutputParser
from langchain_core.prompts import ChatPromptTemplate

# PowerPoint generation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR

from dotenv import load_dotenv
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Configuration
GROQ_API_KEY = os.getenv("GROQ_API_KEY", "")
if not GROQ_API_KEY:
    raise ValueError("GROQ_API_KEY not found in environment variables")

# Enhanced content structure with styling
class ContentBlock(TypedDict):
    type: str  # 'heading', 'paragraph', 'bullet_list', 'numbered_list', 'quote'
    text: str
    style: Dict[str, Any]  # font_size, font_weight, color, etc.
    level: Optional[int]  # for bullet points and headings

class SlideContent(TypedDict):
    title: str
    content_blocks: List[ContentBlock]
    slide_type: str
    layout: str

# State definition for LangGraph
class PresentationState(TypedDict):
    messages: Annotated[List[BaseMessage], add_messages]
    document_content: str
    document_chunks: List[str]
    summary: str
    outline: List[Dict[str, str]]
    slides_content: List[SlideContent]
    session_id: str
    status: str

# Document processor with vectorstore
class DocumentProcessor:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=500,
            chunk_overlap=100,
            length_function=len,
        )
        # Use free HuggingFace embeddings
        self.embeddings = HuggingFaceEndpointEmbeddings(
            model='sentence-transformers/all-MiniLM-L6-v2',
            task="feature-extraction",
            huggingfacehub_api_token=os.getenv('HUGGINGFACEHUB_API_TOKEN'))
        
        self.vectorstore = None
    
    def process_pdf(self, file_path: str) -> tuple[str, List[str]]:
        """Load PDF and create chunks"""
        loader = PyPDFLoader(file_path)
        documents = loader.load()
        
        # Combine all pages
        full_content = "\n".join([doc.page_content for doc in documents])
        
        # Create chunks
        chunks = self.text_splitter.split_text(full_content)
        
        # Create vectorstore
        self.vectorstore = FAISS.from_texts(chunks, self.embeddings)
        
        logger.info(f"PDF processed: {len(chunks)} chunks created")
        return full_content, chunks
    
    def retrieve_relevant_chunks(self, query: str, k: int = 3) -> List[str]:
        """Retrieve relevant chunks using similarity search"""
        if not self.vectorstore:
            return []
        
        docs = self.vectorstore.similarity_search(query, k=k)
        return [doc.page_content for doc in docs]

# Enhanced PowerPoint generator with styling support
class EnhancedPowerPointGenerator:
    def __init__(self):
        self.presentation = None
        # Define color palette
        self.colors = {
            'primary': RGBColor(31, 73, 125),      # Professional blue
            'secondary': RGBColor(68, 114, 196),   # Light blue
            'accent': RGBColor(112, 173, 71),      # Green
            'text_dark': RGBColor(64, 64, 64),     # Dark gray
            'text_light': RGBColor(89, 89, 89),    # Light gray
            'background': RGBColor(248, 249, 250)  # Very light gray
        }
    
    def create_presentation(self, slides_data: List[SlideContent], title: str = "AI Generated Presentation") -> BytesIO:
        """Create PowerPoint presentation from enhanced slides data"""
        # Create new presentation
        prs = Presentation()
        
        # Title slide
        self._create_title_slide(prs, title)
        
        # Content slides
        for slide_data in slides_data:
            self._add_enhanced_content_slide(prs, slide_data)
        
        # Save presentation
        ppt_stream = BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)
        
        logger.info(f"Enhanced PowerPoint created with {len(slides_data) + 1} slides")
        return ppt_stream
    
    def _create_title_slide(self, prs: Presentation, title: str):
        """Create enhanced title slide"""
        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = title
        title_slide.placeholders[1].text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        
        # Style title slide
        title_shape = title_slide.shapes.title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(44)
            paragraph.font.bold = True
            paragraph.font.color.rgb = self.colors['primary']
            paragraph.alignment = PP_ALIGN.CENTER
        
        subtitle_shape = title_slide.placeholders[1]
        for paragraph in subtitle_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = self.colors['text_light']
            paragraph.alignment = PP_ALIGN.CENTER
    
    def _add_enhanced_content_slide(self, prs: Presentation, slide_data: SlideContent):
        """Add enhanced content slide with proper styling"""
        # Determine layout based on slide type and content
        layout_type = self._determine_layout(slide_data)
        slide_layout = prs.slide_layouts[layout_type]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = slide_data["title"]
        self._style_title(title)
        
        # Add content blocks
        if layout_type == 1:  # Title and Content layout
            content_placeholder = slide.placeholders[1]
            self._populate_content_placeholder(content_placeholder, slide_data["content_blocks"])
        else:  # Other layouts - create text boxes
            self._add_content_textboxes(slide, slide_data["content_blocks"])
    
    def _determine_layout(self, slide_data: SlideContent) -> int:
        """Determine best layout based on content"""
        content_blocks = slide_data["content_blocks"]
        
        # Simple heuristic for now
        if len(content_blocks) <= 4 and all(block["type"] in ["bullet_list", "paragraph"] for block in content_blocks):
            return 1  # Title and Content
        elif slide_data["slide_type"] == "intro":
            return 0  # Title slide layout
        else:
            return 1  # Default to title and content
    
    def _style_title(self, title_shape):
        """Apply consistent title styling"""
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(36)
            paragraph.font.bold = True
            paragraph.font.color.rgb = self.colors['primary']
            paragraph.alignment = PP_ALIGN.LEFT
    
    def _populate_content_placeholder(self, content_placeholder, content_blocks: List[ContentBlock]):
        """Populate content placeholder with styled content blocks"""
        content_placeholder.text = ""
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        first_paragraph = True
        
        for block in content_blocks:
            if block["type"] == "heading":
                p = text_frame.add_paragraph() if not first_paragraph else text_frame.paragraphs[0]
                p.text = block["text"]
                self._apply_heading_style(p, block)
                
            elif block["type"] == "paragraph":
                p = text_frame.add_paragraph() if not first_paragraph else text_frame.paragraphs[0]
                p.text = block["text"]
                self._apply_paragraph_style(p, block)
                
            elif block["type"] == "bullet_list":
                items = self._parse_list_items(str(block["text"]))
                for i, item in enumerate(items):
                    p = text_frame.add_paragraph() if not first_paragraph or i > 0 else text_frame.paragraphs[0]
                    p.text = item
                    p.level = block.get("level", 0)
                    self._apply_bullet_style(p, block)
                    
            elif block["type"] == "numbered_list":
                items = self._parse_list_items(str(block["text"]))
                for i, item in enumerate(items):
                    p = text_frame.add_paragraph() if not first_paragraph or i > 0 else text_frame.paragraphs[0]
                    p.text = f"{i+1}. {item}"
                    self._apply_numbered_style(p, block)
                    
            elif block["type"] == "quote":
                p = text_frame.add_paragraph() if not first_paragraph else text_frame.paragraphs[0]
                p.text = f'"{block["text"]}"'
                self._apply_quote_style(p, block)
            
            first_paragraph = False
    
    def _parse_list_items(self, text: str) -> List[str]:
        """Parse list items from text, handling various bullet formats"""
        # Remove common bullet characters and clean up
        text = text.replace("•", "\n").replace("*", "\n").replace("-", "\n")
        text = text.replace("Ã¢â‚¬Â¢", "\n")  # Fix encoding issues
        
        items = []
        for line in text.split('\n'):
            line = line.strip()
            if line and len(line) > 3:  # Avoid very short lines
                # Remove leading numbers or bullets
                if line[0].isdigit() and '. ' in line[:5]:
                    line = line.split('. ', 1)[1]
                items.append(line)
        
        return items[:6]  # Limit to 6 items per slide
    
    def _apply_heading_style(self, paragraph, block: ContentBlock):
        """Apply heading style"""
        style = block.get("style", {})
        paragraph.font.size = Pt(style.get("font_size", 24))
        paragraph.font.bold = style.get("font_weight", "bold") == "bold"
        paragraph.font.color.rgb = self.colors.get(style.get("color", "primary"), self.colors['primary'])
        paragraph.space_after = Pt(12)
    
    def _apply_paragraph_style(self, paragraph, block: ContentBlock):
        """Apply paragraph style"""
        style = block.get("style", {})
        paragraph.font.size = Pt(style.get("font_size", 16))
        paragraph.font.bold = style.get("font_weight", "normal") == "bold"
        paragraph.font.color.rgb = self.colors.get(style.get("color", "text_dark"), self.colors['text_dark'])
        paragraph.space_after = Pt(8)
        paragraph.alignment = PP_ALIGN.LEFT
    
    def _apply_bullet_style(self, paragraph, block: ContentBlock):
        """Apply bullet point style"""
        style = block.get("style", {})
        paragraph.font.size = Pt(style.get("font_size", 14))
        paragraph.font.bold = style.get("font_weight", "normal") == "bold"
        paragraph.font.color.rgb = self.colors.get(style.get("color", "text_dark"), self.colors['text_dark'])
        paragraph.space_after = Pt(6)
        paragraph.alignment = PP_ALIGN.LEFT
    
    def _apply_numbered_style(self, paragraph, block: ContentBlock):
        """Apply numbered list style"""
        style = block.get("style", {})
        paragraph.font.size = Pt(style.get("font_size", 14))
        paragraph.font.bold = style.get("font_weight", "normal") == "bold"
        paragraph.font.color.rgb = self.colors.get(style.get("color", "text_dark"), self.colors['text_dark'])
        paragraph.space_after = Pt(6)
    
    def _apply_quote_style(self, paragraph, block: ContentBlock):
        """Apply quote style"""
        style = block.get("style", {})
        paragraph.font.size = Pt(style.get("font_size", 16))
        paragraph.font.italic = True
        paragraph.font.color.rgb = self.colors.get(style.get("color", "secondary"), self.colors['secondary'])
        paragraph.space_after = Pt(10)
        paragraph.alignment = PP_ALIGN.CENTER
    
    def _add_content_textboxes(self, slide, content_blocks: List[ContentBlock]):
        """Add content as separate text boxes for more complex layouts"""
        # This method can be expanded for more complex slide layouts
        # For now, fall back to the main content approach
        pass

# Enhanced LangGraph workflow
class EnhancedPresentationWorkflow:
    def __init__(self):
        # Use different models for different tasks
        self.fast_llm = ChatGroq(
            model="llama-3.1-8b-instant",
            groq_api_key=GROQ_API_KEY,
            temperature=0.3
        )
        self.smart_llm = ChatGroq(
            model="llama-3.3-70b-versatile", 
            groq_api_key=GROQ_API_KEY,
            temperature=0.5
        )
        
        self.doc_processor = DocumentProcessor()
        self.ppt_generator = EnhancedPowerPointGenerator()
        self.graph = self._build_graph()
    
    def _build_graph(self) -> StateGraph:
        """Build the presentation generation workflow"""
        workflow = StateGraph(PresentationState)
        
        # Add nodes
        workflow.add_node("summarize", self.summarize_document)
        workflow.add_node("create_outline", self.create_presentation_outline)  
        workflow.add_node("generate_content", self.generate_enhanced_slide_content)
        
        # Define flow
        workflow.set_entry_point("summarize")
        workflow.add_edge("summarize", "create_outline")
        workflow.add_edge("create_outline", "generate_content")
        workflow.add_edge("generate_content", END)
        
        return workflow.compile()
    
    def summarize_document(self, state: PresentationState) -> PresentationState:
        """Summarize the entire document using smart LLM with chunking for large docs"""
        try:
            content = state["document_content"]
            
            # For very large documents, create a summary from chunks first
            if len(content) > 40000:  # ~10k tokens
                logger.info("Large document detected, using chunk-based summarization")
                chunks = state["document_chunks"][:10]  # Use first 10 chunks
                
                chunk_summaries = []
                for i, chunk in enumerate(chunks):
                    try:
                        chunk_prompt = ChatPromptTemplate.from_template(
                            """Summarize this section of a document in 2-3 sentences, focusing on key points:

{chunk}

Summary:"""
                        )
                        
                        response = self.fast_llm.invoke(
                            chunk_prompt.format(chunk=chunk)
                        )
                        chunk_summaries.append(response.content.strip())
                        logger.info(f"Processed chunk {i+1}/{len(chunks)}")
                    except Exception as e:
                        logger.warning(f"Failed to process chunk {i+1}: {e}")
                        continue
                
                # Combine chunk summaries
                combined_summary = "\n\n".join(chunk_summaries)
                
                # Final comprehensive summary
                final_prompt = ChatPromptTemplate.from_template(
                    """Based on these section summaries, create a comprehensive document summary:

{summaries}

Create a detailed summary (300-500 words) covering:
1. Main topic and purpose
2. Key findings or arguments  
3. Important data points or evidence
4. Conclusions or recommendations"""
                )
                
                final_response = self.smart_llm.invoke(
                    final_prompt.format(summaries=combined_summary)
                )
                
                state["summary"] = final_response.content
                
            else:
                # Use original approach for smaller documents
                summary_prompt = ChatPromptTemplate.from_template(
                    """Analyze the following document and create a comprehensive summary that captures:
1. Main topic and purpose
2. Key findings or arguments
3. Important data points or evidence
4. Conclusions or recommendations

Document:
{content}

Provide a detailed summary in 300-500 words that will serve as the foundation for creating a presentation."""
                )
                
                # Truncate content if too long (roughly 8000 tokens for safety)
                truncated_content = content[:32000] if len(content) > 32000 else content
                
                response = self.smart_llm.invoke(
                    summary_prompt.format(content=truncated_content)
                )
                
                state["summary"] = response.content
            
            state["status"] = "summarized"
            logger.info("Document summarized successfully")
            
            return state
            
        except Exception as e:
            logger.error(f"Error in summarization: {e}")
            state["status"] = "error"
            return state
    
    def create_presentation_outline(self, state: PresentationState) -> PresentationState:
        """Create presentation structure using smart LLM with exactly 10 slides"""
        try:
            summary = state["summary"]
            
            outline_prompt = ChatPromptTemplate.from_template(
                """Based on this document summary, create a presentation outline with EXACTLY 10 slides.

Summary:
{summary}

Return ONLY a JSON array with this exact format:
[
  {{"title": "Introduction", "description": "Overview of the topic", "slide_type": "intro"}},
  {{"title": "Background", "description": "Context and background", "slide_type": "content"}},
  {{"title": "Key Findings", "description": "Main discoveries", "slide_type": "content"}},
  {{"title": "Conclusion", "description": "Summary and next steps", "slide_type": "conclusion"}}
]

IMPORTANT: Create exactly 10 slides total:
- 1 intro slide
- 8 content slides covering main points  
- 1 conclusion slide

Each title should be concise (2-6 words). Descriptions should be brief (5-10 words)."""
            )
            
            try:
                # Use smart LLM for outline generation as requested
                response = self.smart_llm.invoke(
                    outline_prompt.format(summary=summary)
                )
                
                # Parse JSON response
                parser = JsonOutputParser()
                outline = parser.parse(response.content)
                
                # Validate and ensure we have exactly 10 slides
                if not isinstance(outline, list):
                    raise ValueError("Invalid outline format")
                    
                # Ensure exactly 10 slides
                if len(outline) > 10:
                    outline = outline[:10]  # Take first 10
                elif len(outline) < 10:
                    # If less than 10, use fallback to ensure we have exactly 10
                    raise ValueError("Insufficient slides generated")
                    
                state["outline"] = outline
                logger.info(f"Outline created with exactly {len(state['outline'])} slides")
                
            except Exception as parse_error:
                logger.warning(f"Outline generation failed: {parse_error}, using fallback")
                # Smart fallback based on document type - exactly 10 slides

                state["outline"] = [
                        {"title": "Introduction", "description": "Document overview", "slide_type": "intro"},
                        {"title": "Background", "description": "Context and setting", "slide_type": "content"}, 
                        {"title": "Main Concepts", "description": "Key ideas", "slide_type": "content"},
                        {"title": "Key Findings", "description": "Important discoveries", "slide_type": "content"},
                        {"title": "Analysis", "description": "Detailed examination", "slide_type": "content"},
                        {"title": "Evidence", "description": "Supporting data", "slide_type": "content"},
                        {"title": "Implications", "description": "What this means", "slide_type": "content"},
                        {"title": "Applications", "description": "Practical uses", "slide_type": "content"},
                        {"title": "Future Directions", "description": "Next steps", "slide_type": "content"},
                        {"title": "Conclusion", "description": "Summary and takeaways", "slide_type": "conclusion"}
                    ]
            
            # Ensure we have exactly 10 slides
            assert len(state["outline"]) == 10, f"Expected 10 slides, got {len(state['outline'])}"
            
            state["status"] = "outlined"
            
            return state
            
        except Exception as e:
            logger.error(f"Error creating outline: {e}")
            state["status"] = "error"
            return state
    
    def generate_enhanced_slide_content(self, state: PresentationState) -> PresentationState:
        """Generate enhanced slide content with styling information"""
        try:
            outline = state["outline"]
            
            # Ensure we have exactly 10 slides
            if len(outline) != 10:
                logger.error(f"Expected exactly 10 slides in outline, got {len(outline)}")
                state["status"] = "error"
                return state
            
            slides_content = []
            
            enhanced_content_prompt = ChatPromptTemplate.from_template(
                """Create structured slide content for this presentation slide:

Title: {title}
Description: {description}
Slide Type: {slide_type}
Relevant Information: {context}

Generate a JSON response with the following structure:
{{
  "title": "{title}",
  "content_blocks": [
    {{
      "type": "paragraph|bullet_list|heading|quote",
      "text": "content text",
      "style": {{
        "font_size": 16,
        "font_weight": "normal|bold",
        "color": "text_dark|primary|secondary",
        "alignment": "left|center"
      }},
      "level": 0
    }}
  ],
  "slide_type": "{slide_type}",
  "layout": "title_content"
}}

Content guidelines:
- Use mix of paragraphs, bullet lists, and headings for variety
- Intro slides: Use larger headings and introductory paragraphs
- Content slides: Mix headings, paragraphs, and bullet points
- Conclusion slides: Use summary paragraphs and key takeaways
- Keep bullet lists to 4-6 items maximum
- Make content informative but readable for presentations"""
            )
            
            # Process each slide individually - exactly 10 LLM calls
            for i, slide_info in enumerate(outline, 1):
                title = slide_info["title"]
                description = slide_info["description"]
                slide_type = slide_info.get("slide_type", "content")
                
                logger.info(f"Processing enhanced slide {i}/10: {title}")
                
                try:
                    # Retrieve relevant chunks for this slide topic
                    query = f"{title} {description}"
                    relevant_chunks = self.doc_processor.retrieve_relevant_chunks(query, k=2)
                    
                    if relevant_chunks:
                        context = "\n".join(relevant_chunks)
                    else:
                        # Fallback to summary if no chunks found
                        context = state["summary"][:1000]  # Limit context size
                    
                    # Generate enhanced content using smart LLM for better structure
                    response = self.smart_llm.invoke(
                        enhanced_content_prompt.format(
                            title=title,
                            description=description,
                            slide_type=slide_type,
                            context=context
                        )
                    )
                    
                    try:
                        # Parse JSON response
                        parser = JsonOutputParser()
                        slide_content = parser.parse(response.content)
                        
                        # Validate structure
                        if not isinstance(slide_content, dict) or "content_blocks" not in slide_content:
                            raise ValueError("Invalid slide content structure")
                        
                        slides_content.append(slide_content)
                        logger.info(f"Successfully generated enhanced content for slide {i}: {title}")
                        
                    except Exception as json_error:
                        logger.warning(f"JSON parsing failed for slide {i}: {json_error}, using fallback")
                        # Create fallback structured content
                        fallback_content = self._generate_fallback_structured_content(
                            title, description, slide_type, context
                        )
                        slides_content.append(fallback_content)
                        
                except Exception as e:
                    logger.warning(f"Failed to generate content for slide {i} '{title}': {e}")
                    # Add fallback content
                    fallback_content = self._generate_fallback_structured_content(
                        title, description, slide_type, state["summary"][:500]
                    )
                    slides_content.append(fallback_content)
                    logger.info(f"Used fallback structured content for slide {i}: {title}")
            
            # Ensure we generated exactly 10 slides
            if len(slides_content) != 10:
                logger.error(f"Expected 10 slides, generated {len(slides_content)}")
                state["status"] = "error"
                return state
            
            state["slides_content"] = slides_content
            state["status"] = "completed"
            logger.info(f"Successfully generated enhanced content for all 10 slides")
            
            return state
            
        except Exception as e:
            logger.error(f"Error generating enhanced slide content: {e}")
            state["status"] = "error"
            return state
    
    def _generate_fallback_structured_content(self, title: str, description: str, slide_type: str, context: str) -> SlideContent:
        """Generate fallback structured content when JSON parsing fails"""
        
        # Extract key sentences from context
        sentences = [s.strip() for s in context.split('.') if len(s.strip()) > 20]
        
        if slide_type == "intro":
            content_blocks = [
                {
                    "type": "heading",
                    "text": description.title(),
                    "style": {"font_size": 24, "font_weight": "bold", "color": "primary"},
                    "level": 1
                },
                {
                    "type": "paragraph", 
                    "text": sentences[0] + "." if sentences else "This presentation covers the key aspects of the topic.",
                    "style": {"font_size": 18, "font_weight": "normal", "color": "text_dark"}
                }
            ]
        elif slide_type == "conclusion":
            content_blocks = [
                {
                    "type": "heading",
                    "text": "Key Takeaways",
                    "style": {"font_size": 20, "font_weight": "bold", "color": "primary"},
                    "level": 1
                },
                {
                    "type": "bullet_list",
                    "text": "\n".join(sentences[:4]) if len(sentences) >= 4 else "Important findings and conclusions\nKey insights from the analysis\nImplications for future work",
                    "style": {"font_size": 16, "font_weight": "normal", "color": "text_dark"},
                    "level": 0
                }
            ]
        else:  # content slide
            title_words = title.lower().split()
            relevant_sentences = [s for s in sentences if any(word in s.lower() for word in title_words)]
            
            if relevant_sentences:
                content_blocks = [
                    {
                        "type": "paragraph",
                        "text": relevant_sentences[0] + ".",
                        "style": {"font_size": 14, "font_weight": "normal", "color": "text_dark"},
                        "level": 0
                    }
                ]
        
        return {
            "title": title,
            "content_blocks": content_blocks,
            "slide_type": slide_type,
            "layout": "title_content"
        }

# Main orchestrator - updated to use enhanced workflow
class PresentationOrchestrator:
    def __init__(self):
        self.workflow = EnhancedPresentationWorkflow()  # Use enhanced workflow
        self.sessions = {}
    
    def process_pdf(self, file_path: str, session_id: str) -> Dict[str, Any]:
        """Process PDF document"""
        try:
            content, chunks = self.workflow.doc_processor.process_pdf(file_path)
            
            self.sessions[session_id] = {
                "content": content,
                "chunks": chunks,
                "processed_at": datetime.now(),
                "file_path": file_path,
                "slides_content": None  # Will be populated after generation
            }
            
            return {
                "status": "success",
                "message": "PDF processed successfully", 
                "content_length": len(content),
                "chunks_count": len(chunks)
            }
            
        except Exception as e:
            logger.error(f"Error processing PDF: {e}")
            raise HTTPException(status_code=400, detail=str(e))
    
    def generate_presentation(self, session_id: str) -> tuple[BytesIO, Dict[str, Any]]:
        """Generate presentation and return PowerPoint file stream"""
        if session_id not in self.sessions:
            raise HTTPException(status_code=404, detail="Session not found")
        
        session_data = self.sessions[session_id]
        
        # Check if slides content already exists (from previous generation)
        if session_data.get("slides_content"):
            slides_content = session_data["slides_content"]
        else:
            # Generate slides content if not exists
            initial_state = PresentationState(
                messages=[],
                document_content=session_data["content"],
                document_chunks=session_data["chunks"], 
                summary="",
                outline=[],
                slides_content=[],
                session_id=session_id,
                status="started"
            )
            
            try:
                logger.info(f"Starting enhanced presentation generation for session {session_id}")
                
                # Run the enhanced workflow
                final_state = self.workflow.graph.invoke(initial_state)
                
                if final_state["status"] != "completed":
                    raise HTTPException(
                        status_code=500,
                        detail=f"Generation failed at status: {final_state['status']}"
                    )
                
                slides_content = final_state["slides_content"]
                # Store for future use
                self.sessions[session_id]["slides_content"] = slides_content
                
            except Exception as e:
                logger.error(f"Error generating enhanced presentation: {e}")
                raise HTTPException(status_code=500, detail=str(e))
        
        # Create PowerPoint stream using enhanced generator
        ppt_stream = self.workflow.ppt_generator.create_presentation(
            slides_content,
            title="AI Generated Presentation"
        )
        
        result_data = {
            "slides": slides_content,
            "metadata": {
                "generated_at": datetime.now().isoformat(),
                "total_slides": len(slides_content),
                "session_id": session_id,
                "enhanced": True  # Flag to indicate enhanced format
            }
        }
        
        logger.info(f"Enhanced presentation generated: {len(slides_content)} slides")
        return ppt_stream, result_data

    def generate_presentation_metadata(self, session_id: str) -> Dict[str, Any]:
        """Generate presentation metadata only (for preview)"""
        if session_id not in self.sessions:
            raise HTTPException(status_code=404, detail="Session not found")
        
        session_data = self.sessions[session_id]
        
        # Initial state
        initial_state = PresentationState(
            messages=[],
            document_content=session_data["content"],
            document_chunks=session_data["chunks"], 
            summary="",
            outline=[],
            slides_content=[],
            session_id=session_id,
            status="started"
        )
        
        try:
            logger.info(f"Starting enhanced presentation generation for session {session_id}")
            
            # Run the enhanced workflow
            final_state = self.workflow.graph.invoke(initial_state)
            
            if final_state["status"] != "completed":
                raise HTTPException(
                    status_code=500,
                    detail=f"Generation failed at status: {final_state['status']}"
                )
            
            # Store the generated slides in session for later PowerPoint creation
            self.sessions[session_id]["slides_content"] = final_state["slides_content"]
            
            result_data = {
                "slides": final_state["slides_content"],
                "metadata": {
                    "generated_at": datetime.now().isoformat(),
                    "total_slides": len(final_state["slides_content"]),
                    "session_id": session_id,
                    "enhanced": True  # Flag to indicate enhanced format
                }
            }
            
            logger.info(f"Enhanced presentation generated: {len(final_state['slides_content'])} slides")
            return result_data
            
        except Exception as e:
            logger.error(f"Error generating enhanced presentation: {e}")
            raise HTTPException(status_code=500, detail=str(e))

# FastAPI setup
app = FastAPI(
    title="Enhanced AI Presentation Orchestrator",
    description="Generate styled PowerPoint presentations from PDF documents with enhanced formatting",
    version="2.0.0"
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize orchestrator
orchestrator = PresentationOrchestrator()

# Pydantic models
class GenerateRequest(BaseModel):
    session_id: str

# API Routes
@app.get("/")
async def root():
    return {
        "message": "Enhanced AI Presentation Orchestrator",
        "version": "2.0.0",
        "features": ["Enhanced styling", "Mixed content types", "Professional formatting"],
        "endpoints": ["/upload", "/generate", "/download_ppt/{session_id}", "/health"]
    }

@app.post("/upload")
async def upload_pdf(file: UploadFile = File(...)):
    """Upload and process PDF"""
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Only PDF files are supported")
    
    # Generate session ID
    session_id = f"session_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
    
    # Save uploaded file temporarily
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
        content = await file.read()
        tmp_file.write(content)
        tmp_file_path = tmp_file.name
    
    try:
        result = orchestrator.process_pdf(tmp_file_path, session_id)
        result["session_id"] = session_id
        return result
        
    except Exception as e:
        # Clean up temp file
        Path(tmp_file_path).unlink(missing_ok=True)
        raise e

@app.post("/generate")
async def generate_presentation(request: GenerateRequest):
    """Generate enhanced presentation metadata (for preview)"""
    return orchestrator.generate_presentation_metadata(request.session_id)

@app.get("/download_ppt/{session_id}")
async def download_presentation(session_id: str):
    """Generate and download enhanced PowerPoint file directly from session"""
    try:
        # Generate presentation and get PowerPoint stream
        ppt_stream, result_data = orchestrator.generate_presentation(session_id)
        
        # Reset stream position
        ppt_stream.seek(0)
        
        # Create filename with timestamp
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"enhanced_presentation_{session_id}_{timestamp}.pptx"
        
        # Return as streaming response
        return StreamingResponse(
            io.BytesIO(ppt_stream.read()),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error downloading enhanced presentation: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/sessions")
async def list_sessions():
    """List all sessions with enhanced info"""
    sessions_info = {}
    for session_id, data in orchestrator.sessions.items():
        sessions_info[session_id] = {
            "processed_at": data["processed_at"].isoformat(),
            "content_length": len(data["content"]),
            "chunks_count": len(data["chunks"]),
            "has_slides": data.get("slides_content") is not None,
            "enhanced": True
        }
    return {"sessions": sessions_info}

@app.get("/health")
async def health_check():
    """Health check with enhanced features info"""
    try:
        test_response = orchestrator.workflow.fast_llm.invoke("test")
        llm_status = "healthy"
    except Exception as e:
        llm_status = f"error: {str(e)}"
    
    return {
        "status": "healthy",
        "timestamp": datetime.now().isoformat(),
        "version": "2.0.0 Enhanced",
        "components": {
            "api": "healthy",
            "llm": llm_status,
            "sessions": len(orchestrator.sessions),
            "enhanced_generator": "active"
        },
        "models": {
            "fast": "llama-3.1-8b-instant",
            "smart": "llama-3.3-70b-versatile"
        },
        "features": [
            "Character encoding fixes",
            "Mixed content types",
            "Enhanced styling",
            "Professional formatting",
            "JSON-structured content"
        ]
    }

# Add new endpoint for slide preview
@app.get("/preview/{session_id}")
async def preview_slides(session_id: str):
    """Preview slide content without generating PowerPoint"""
    if session_id not in orchestrator.sessions:
        raise HTTPException(status_code=404, detail="Session not found")
    
    session_data = orchestrator.sessions[session_id]
    slides_content = session_data.get("slides_content")
    
    if not slides_content:
        raise HTTPException(status_code=404, detail="No slides generated for this session")
    
    # Return formatted preview
    preview_data = {
        "session_id": session_id,
        "total_slides": len(slides_content),
        "slides": []
    }
    
    for i, slide in enumerate(slides_content, 1):
        slide_preview = {
            "slide_number": i,
            "title": slide["title"],
            "slide_type": slide["slide_type"],
            "content_blocks_count": len(slide["content_blocks"]),
            "content_blocks": slide["content_blocks"]
        }
        preview_data["slides"].append(slide_preview)
    
    return preview_data

if __name__ == "__main__":
    logger.info("Starting Enhanced AI Presentation Orchestrator v2.0...")
    uvicorn.run(app, host="0.0.0.0", port=8000)
